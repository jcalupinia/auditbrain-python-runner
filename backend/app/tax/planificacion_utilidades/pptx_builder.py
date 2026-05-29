"""Generador de la presentación ejecutiva (.pptx) — diseño premium.

Estética: dark premium consistente (Midnight Executive adaptada a marca
AuditBrain), tipografía DM Sans, cifras como callouts grandes, layouts variados,
motivo visual de círculos dorados numerados. Sin barras decorativas ni líneas
bajo títulos (anti "AI slop"). Deck 16:9 alimentado por el `content` del frontend.

Compatible con python-pptx 0.6.23 (producción).
"""

from __future__ import annotations

import io

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# Paleta (Midnight Executive · marca AuditBrain)
NAVY = RGBColor(0x0A, 0x23, 0x42)   # fondo dominante
DEEP = RGBColor(0x05, 0x14, 0x26)   # fondo portada/cierre (más oscuro)
GOLD = RGBColor(0xC7, 0xA8, 0x3C)   # acento
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ICE = RGBColor(0xCA, 0xDC, 0xFC)    # azul hielo (texto secundario)
GREY = RGBColor(0x9F, 0xB0, 0xC4)
CARD = RGBColor(0x10, 0x2E, 0x4F)   # tarjetas
LINE = RGBColor(0x21, 0x40, 0x66)
GREEN = RGBColor(0x4C, 0xC2, 0x8A)
RED = RGBColor(0xE0, 0x6A, 0x5C)
BLUE = RGBColor(0x5B, 0x9B, 0xE0)
FONT = "DM Sans"

W = Inches(13.333)
H = Inches(7.5)


def build_deck(content: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    c = content or {}
    _cover(prs, c)
    _resumen(prs, c)
    _historico(prs, c)
    _diag_financiero(prs, c)
    _diag_trib_soc(prs, c)
    _alternativas(prs, c)
    _matriz(prs, c)
    _grafico_escenarios(prs, c)
    _recomendacion(prs, c)
    _modelacion(prs, c)
    _plan(prs, c)
    _cierre(prs, c)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ----------------------------------------------------------------- helpers
def _bg(slide, color=NAVY):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r


def _blank(prs, color=NAVY):
    return _bg(prs.slides.add_slide(prs.slide_layouts[6]), color)


def _txt(slide, x, y, w, h, text, size, *, color=WHITE, bold=False,
         align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP,
         spacing=None, line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else [str(text)]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing is not None:
            p.space_after = Pt(spacing)
        if line_spacing is not None:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def _circle_num(slide, x, y, n, d=Inches(0.5)):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    c.fill.solid(); c.fill.fore_color.rgb = GOLD
    c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.name = FONT; r.font.size = Pt(16); r.font.bold = True
    r.font.color.rgb = DEEP
    return c


# título de sección: número en círculo dorado + título grande (sin líneas/barras)
def _title(slide, num, text, sub=None):
    _circle_num(slide, Inches(0.7), Inches(0.62), num)
    _txt(slide, Inches(1.4), Inches(0.55), Inches(11.3), Inches(0.9), text, 30, bold=True)
    if sub:
        _txt(slide, Inches(1.4), Inches(1.25), Inches(11.3), Inches(0.5), sub, 13, color=ICE)


def _card(slide, x, y, w, h, color=CARD):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.color.rgb = LINE; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = 0.06
    except Exception:
        pass
    return sh


def _stat(slide, x, y, w, h, value, label, *, vcolor=GOLD, vsize=34):
    _card(slide, x, y, w, h)
    _txt(slide, x + Inches(0.25), y + Inches(0.22), w - Inches(0.5), h - Inches(0.9),
         str(value), vsize, bold=True, color=vcolor, anchor=MSO_ANCHOR.MIDDLE)
    _txt(slide, x + Inches(0.25), y + h - Inches(0.55), w - Inches(0.5), Inches(0.45),
         label, 11, color=GREY)


def _money(v):
    try:
        return "$" + format(int(round(float(v))), ",d")
    except (TypeError, ValueError):
        return str(v)


# ------------------------------------------------------------------ slides
def _cover(prs, c):
    _blank(prs, DEEP)
    s = prs.slides[-1]
    # motivo: bloque lateral dorado (no full-width)
    side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.28), H)
    side.fill.solid(); side.fill.fore_color.rgb = GOLD
    side.line.fill.background(); side.shadow.inherit = False
    _txt(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.5),
         "AUDITBRAIN  ·  EXECUTIVE ADVISORY  ·  TAX ADVISORY", 12, color=GOLD, bold=True)
    _txt(s, Inches(0.85), Inches(2.3), Inches(11.6), Inches(2.2),
         "Planificación tributaria\nsobre utilidades no distribuidas", 38, bold=True,
         line_spacing=1.05)
    _txt(s, Inches(0.9), Inches(4.85), Inches(11.5), Inches(0.6),
         c.get("empresa", "la Compañía"), 24, color=GOLD, bold=True)
    meta = []
    if c.get("ruc"):
        meta.append("RUC " + str(c["ruc"]))
    if c.get("representante"):
        meta.append(str(c["representante"]) + "  ·  Representante legal")
    linea = "Horizonte 2026–2028"
    if c.get("fecha_analisis"):
        linea += "   ·   " + str(c["fecha_analisis"])
    if c.get("fecha_corte"):
        linea += "   ·   Corte: " + str(c["fecha_corte"])
    meta.append(linea)
    _txt(s, Inches(0.9), Inches(5.5), Inches(11.5), Inches(1.2),
         "\n".join(meta), 14, color=ICE, line_spacing=1.3)
    _txt(s, Inches(0.9), Inches(6.95), Inches(11.5), Inches(0.4),
         "Confidencial · documento preliminar sujeto a revisión y aprobación", 10,
         color=GREY, italic=True)


def _resumen(prs, c):
    prs_slide = prs.slides[-1] if False else None
    _blank(prs)
    s = prs.slides[-1]
    _title(s, 1, "Resumen Ejecutivo")
    rec = (c.get("recomendacion") or "")[:160]
    if rec:
        _txt(s, Inches(1.4), Inches(1.2), Inches(11.3), Inches(0.6), rec, 13, color=ICE)
    kpis = c.get("kpis", [])[:6]
    cols, x0, y0 = 3, Inches(0.7), Inches(2.05)
    cw, ch, gx, gy = Inches(3.95), Inches(2.0), Inches(0.18), Inches(0.28)
    palette = [GOLD, RED, GREEN, GOLD, BLUE, RED]
    for i, k in enumerate(kpis):
        rr, col = divmod(i, cols)
        x = x0 + col * (cw + gx)
        y = y0 + rr * (ch + gy)
        _stat(s, x, y, cw, ch, k.get("valor", ""), k.get("label", ""),
              vcolor=palette[i % len(palette)], vsize=32)


def _table(slide, x, y, w, headers, rows, col_w=None, highlight_row=None, left_cols=(0,)):
    nrows, ncols = len(rows) + 1, len(headers)
    h = Inches(0.46) * nrows
    table = slide.shapes.add_table(nrows, ncols, x, y, w, h).table
    table.first_row = False
    table.horz_banding = False
    if col_w:
        for ci, cw in enumerate(col_w):
            table.columns[ci].width = cw
    for ci, htxt in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = DEEP
        _cell(cell, htxt, 11, GOLD, bold=True,
              align=PP_ALIGN.LEFT if ci in left_cols else PP_ALIGN.RIGHT)
    for ri, row in enumerate(rows, start=1):
        hl = highlight_row is not None and ri - 1 == highlight_row
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD if hl else NAVY
            _cell(cell, str(val), 11, GOLD if hl else WHITE, bold=hl,
                  align=PP_ALIGN.LEFT if ci in left_cols else PP_ALIGN.RIGHT)
    return table


def _cell(cell, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.12)
    cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color


def _historico(prs, c):
    _blank(prs)
    s = prs.slides[-1]
    _title(s, 2, "Pago Históricamente Realizado",
           "Base = utilidades retenidas del cierre del ejercicio anterior")
    headers = ["Año", "Base (acum. anterior)", "Pago a cuenta", "Estado"]
    rows = []
    for h in c.get("pago_historico", []):
        rows.append([h.get("anio"), _money(h.get("base")), _money(h.get("pago")),
                     h.get("estado", "Realizado")])
    mod = c.get("modelacion_2026_2028", {})
    anios = mod.get("anios", []); pagos = mod.get("pago_a_cuenta", [])
    for i, a in enumerate(anios):
        rows.append([a, "—", _money(pagos[i] if i < len(pagos) else 0), "Proyectado"])
    _table(s, Inches(1.4), Inches(2.1), Inches(10.5), headers, rows,
           col_w=[Inches(1.7), Inches(3.9), Inches(2.9), Inches(2.0)])


def _diag_financiero(prs, c):
    _blank(prs)
    s = prs.slides[-1]
    _title(s, 3, "Diagnóstico Financiero")
    d = c.get("diagnostico_financiero", {})
    items = [
        ("Liquidez corriente", d.get("liquidez", "—"), GREEN),
        ("Endeudamiento", d.get("endeudamiento", "—"), GOLD),
        ("ROE", d.get("roe", "—"), BLUE),
        ("Margen neto", d.get("margen_neto", "—"), GREEN),
        ("Días de inventario", d.get("dias_inventario", "—"), RED),
        ("Ciclo de efectivo", d.get("ciclo_efectivo", "—"), RED),
    ]
    cols, x0, y0 = 3, Inches(0.7), Inches(2.05)
    cw, ch, gx, gy = Inches(3.95), Inches(2.0), Inches(0.18), Inches(0.28)
    for i, (lab, val, col) in enumerate(items):
        rr, cc = divmod(i, cols)
        _stat(s, x0 + cc * (cw + gx), y0 + rr * (ch + gy), cw, ch, val, lab, vcolor=col, vsize=30)


def _icon_rows(slide, x, y, w, items, gap=Inches(1.15)):
    for i, it in enumerate(items):
        yy = y + i * gap
        _circle_num(slide, x, yy, i + 1, d=Inches(0.6))
        _txt(slide, x + Inches(0.9), yy + Inches(0.02), w - Inches(1.0), Inches(1.0),
             it, 16, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)


def _diag_trib_soc(prs, c):
    _blank(prs)
    s = prs.slides[-1]
    _title(s, 4, "Diagnóstico Tributario y Societario")
    t = c.get("diagnostico_tributario", {})
    soc = c.get("diagnostico_societario", {})
    # dos columnas de stats
    _txt(s, Inches(0.7), Inches(2.0), Inches(5.9), Inches(0.5), "TRIBUTARIO", 14, bold=True, color=GOLD)
    tr = [("Base año 1", t.get("base_anio1", "—")), ("Tarifa", t.get("tarifa", "—")),
          ("Pago año 1", t.get("pago_anio1", "—")), ("Sin acción 2026–28", t.get("pago_horizonte", "—"))]
    for i, (lab, val) in enumerate(tr):
        rr, cc = divmod(i, 2)
        _stat(s, Inches(0.7) + cc * Inches(3.05), Inches(2.5) + rr * Inches(1.95),
              Inches(2.85), Inches(1.7), val, lab, vcolor=ICE, vsize=22)
    _txt(s, Inches(6.95), Inches(2.0), Inches(5.9), Inches(0.5), "SOCIETARIO", 14, bold=True, color=GOLD)
    so = [("Capital", soc.get("capital", "—")), ("Reservas", soc.get("reservas", "—")),
          ("Result. acumulados", soc.get("resultados_acumulados", "—")),
          ("Patrimonio total", soc.get("patrimonio_total", "—"))]
    for i, (lab, val) in enumerate(so):
        rr, cc = divmod(i, 2)
        _stat(s, Inches(6.95) + cc * Inches(3.05), Inches(2.5) + rr * Inches(1.95),
              Inches(2.85), Inches(1.7), val, lab, vcolor=ICE, vsize=22)


def _alternativas(prs, c):
    _blank(prs)
    s = prs.slides[-1]
    _title(s, 5, "Alternativas Previas al 31 de Julio")
    _icon_rows(s, Inches(0.9), Inches(2.2), Inches(11.4), c.get("alternativas", [])[:4])


def _matriz(prs, c):
    _blank(prs)
    s = prs.slides[-1]
    _title(s, 6, "Matriz de Decisión Estratégica", "Decisión aplicada en 2026")
    headers = ["Escenario", "Pago 2026", "Pago 2026–28", "Retención (crédito)",
               "Costo muerto", "Patrimonio 2028"]
    rows, hl = [], None
    for i, m in enumerate(c.get("matriz_escenarios", [])):
        if m.get("recomendado"):
            hl = i
        rows.append([m.get("escenario", ""), _money(m.get("pago_2026")),
                     _money(m.get("pago_2026_2028")), _money(m.get("retencion_credito")),
                     _money(m.get("costo_muerto")), _money(m.get("patrimonio_2028"))])
    _table(s, Inches(0.6), Inches(2.15), Inches(12.1), headers, rows,
           col_w=[Inches(3.9), Inches(1.5), Inches(1.7), Inches(1.9), Inches(1.55), Inches(1.55)],
           highlight_row=hl)
    _txt(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.6),
         "Escenario sugerido resaltado: anula el pago de 2026 y genera crédito recuperable. "
         "La capitalización debe respaldarse con sustancia económica (no inventarios).",
         11, color=GREY, italic=True)


def _add_bar(slide, x, y, w, h, categories, series, colors=None):
    cd = CategoryChartData()
    cd.categories = categories
    for name, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, cd)
    chart = gf.chart
    chart.has_title = False
    try:
        if len(series) > 1:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.color.rgb = WHITE
            chart.legend.font.size = Pt(10); chart.legend.font.name = FONT
        else:
            chart.has_legend = False
        plot = chart.plots[0]
        plot.gap_width = 55
        ca = chart.category_axis
        ca.tick_labels.font.color.rgb = WHITE
        ca.tick_labels.font.size = Pt(11); ca.tick_labels.font.name = FONT
        va = chart.value_axis
        va.tick_labels.font.color.rgb = GREY
        va.tick_labels.font.size = Pt(9); va.has_major_gridlines = False
    except Exception:
        plot = chart.plots[0]
    try:
        if colors and len(series) == 1:
            for i, pt in enumerate(plot.series[0].points):
                pt.format.fill.solid(); pt.format.fill.fore_color.rgb = colors[i % len(colors)]
        else:
            pal = [BLUE, GOLD, GREEN, RED]
            for i, ser in enumerate(plot.series):
                ser.format.fill.solid(); ser.format.fill.fore_color.rgb = pal[i % len(pal)]
    except Exception:
        pass
    return gf


def _grafico_escenarios(prs, c):
    _blank(prs)
    s = prs.slides[-1]
    _title(s, 7, "Pago a Cuenta por Escenario", "Comparativo 2026–2028")
    g = c.get("grafico_pago_por_escenario", {})
    labels = [str(l).split("·")[-1].strip()[:20] for l in g.get("labels", [])]
    vals = g.get("valores", [])
    if labels and vals:
        _add_bar(s, Inches(1.2), Inches(2.1), Inches(10.9), Inches(4.7),
                 labels, [("Pago a cuenta", vals)], colors=[RED, BLUE, GREEN, GOLD])


def _recomendacion(prs, c):
    _blank(prs, DEEP)
    s = prs.slides[-1]
    _title(s, 8, "Recomendación")
    # callout grande del ahorro si está en kpis
    ahorro = next((k.get("valor") for k in c.get("kpis", []) if "horro" in k.get("label", "")), None)
    if ahorro:
        _txt(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(1.3), ahorro, 60, bold=True, color=GOLD)
        _txt(s, Inches(0.95), Inches(3.15), Inches(11.5), Inches(0.5),
             "Ahorro / diferimiento estimado frente a no actuar", 14, color=ICE)
    _card(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(2.2))
    _txt(s, Inches(1.2), Inches(4.15), Inches(10.9), Inches(1.7),
         c.get("recomendacion", ""), 16, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    _txt(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.7),
         c.get("nota", ""), 10, color=GREY, italic=True)


def _modelacion(prs, c):
    _blank(prs)
    s = prs.slides[-1]
    _title(s, 9, "Modelación Financiera y Tributaria", "Proyección 2026–2028")
    mod = c.get("modelacion_2026_2028", {})
    anios = [str(a) for a in mod.get("anios", [])]
    if anios:
        _add_bar(s, Inches(1.2), Inches(2.1), Inches(10.9), Inches(4.7), anios, [
            ("Pago a cuenta", mod.get("pago_a_cuenta", [])),
            ("Crédito vs. IR", mod.get("credito_vs_ir", [])),
            ("En riesgo", mod.get("en_riesgo", [])),
        ])


def _plan(prs, c):
    _blank(prs)
    s = prs.slides[-1]
    _title(s, 10, "Plan de Acción")
    headers = ["Acción", "Responsable", "Plazo"]
    rows = [[p.get("accion", ""), p.get("responsable", ""), p.get("plazo", "")]
            for p in c.get("plan_accion", [])]
    _table(s, Inches(0.7), Inches(2.1), Inches(11.9), headers, rows,
           col_w=[Inches(7.0), Inches(2.7), Inches(2.2)], left_cols=(0, 1, 2))


def _cierre(prs, c):
    _blank(prs, DEEP)
    s = prs.slides[-1]
    side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, W - Inches(0.28), 0, Inches(0.28), H)
    side.fill.solid(); side.fill.fore_color.rgb = GOLD
    side.line.fill.background(); side.shadow.inherit = False
    _txt(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.2), "Gracias", 54, bold=True, color=GOLD)
    _txt(s, Inches(0.95), Inches(3.95), Inches(11.5), Inches(0.6),
         c.get("empresa", ""), 22, bold=True)
    _txt(s, Inches(0.95), Inches(4.7), Inches(11.5), Inches(1.4),
         "AuditBrain®  ·  Tax › Análisis\nProyección de planificación, no auditada. "
         "Régimen sujeto a criterios administrativos del SRI.\nLas cifras normativas "
         "requieren validación profesional.", 12, color=ICE, line_spacing=1.3)
